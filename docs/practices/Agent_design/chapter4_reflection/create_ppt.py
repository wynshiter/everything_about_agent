"""
Chapter 4: Reflection PPT Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_reflection_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    TITLE_COLOR = RGBColor(31, 78, 121)
    ACCENT_COLOR = RGBColor(0, 120, 212)
    
    def add_title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title: str, bullets: list):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.space_after = Pt(12)
    
    # Slide 1: 封面
    add_title_slide("Chapter 4: Reflection", "反思模式 - 自我评估与迭代改进")
    
    # Slide 2: 什么是反思
    add_content_slide("什么是反思模式？", [
        "Agent 评估自己的工作成果并利用评估结果改进性能",
        "引入反馈循环：执行 -> 评估 -> 改进 -> 迭代",
        "最有效的实现：生产者-批评者 (Producer-Critic) 模型",
        "批评者可以是另一个 Agent 或不同的 LLM 调用"
    ])
    
    # Slide 3: 工作流程
    add_content_slide("反思模式的工作流程", [
        "1. 执行：Agent 生成初始输出",
        "2. 评估/批评：分析结果，检查准确性、完整性、风格",
        "3. 反思/改进：根据批评生成改进版本",
        "4. 迭代：重复直到满足条件或达到最大次数"
    ])
    
    # Slide 4: 适用场景
    add_content_slide("反思模式的典型应用", [
        "创意写作：改进文章、故事、营销文案",
        "代码生成与调试：识别错误、修复问题",
        "复杂问题解决：评估中间步骤、排除矛盾",
        "摘要与信息合成：提高准确性",
        "规划与策略：评估计划可行性"
    ])
    
    # Slide 5: LangChain 实现
    add_content_slide("LangChain 实现", [
        "使用循环实现迭代反思",
        "generate() 函数：生成或改进输出",
        "reflect() 函数：评估输出并提供批评",
        "检查 'CODE_IS_PERFECT' 作为停止条件"
    ])
    
    # Slide 6: Google ADK 实现
    add_content_slide("Google ADK 实现", [
        "使用 SequentialAgent 串联生成者和批评者",
        "Generator Agent: output_key='draft'",
        "Reviewer Agent: 读取 state['draft'], 输出到 state['review']",
        "可使用 LoopAgent 实现迭代循环"
    ])
    
    # Slide 7: 运行指南
    add_content_slide("运行指南", [
        "确保已安装依赖：pip install -e .",
        "确保 Ollama 或其他后端服务正在运行",
        "运行示例：python src/agents/patterns/reflection.py"
    ])
    
    # Slide 8: 总结
    add_content_slide("总结与要点", [
        "反思模式显著提高输出质量、准确性和一致性",
        "生产者-批评者模型是最有效的实现方式",
        "代价：增加延迟和计算成本",
        "注意管理上下文长度和迭代历史"
    ])
    
    output_path = "d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter4_reflection/Chapter4_Reflection.pptx"
    prs.save(output_path)
    print(f"Chapter 4 PPT created: {output_path}")

if __name__ == "__main__":
    create_reflection_ppt()
