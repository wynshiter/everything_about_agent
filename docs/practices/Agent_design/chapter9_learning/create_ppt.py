from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    TC = RGBColor(31, 78, 121)
    
    def t(s, s2=""):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold = s, Pt(44), True
        p.font.color.rgb, p.alignment = TC, PP_ALIGN.CENTER
        if s2:
            sb = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            p = sb.text_frame.paragraphs[0]
            p.text, p.font.size = s2, Pt(24)
            p.alignment = PP_ALIGN.CENTER
    
    def c(t, bs):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold = t, Pt(32), True
        p.font.color.rgb = TC
        cb = sl.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = cb.text_frame
        for i, b in enumerate(bs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text, p.font.size = "• " + b, Pt(20)
            p.space_after = Pt(12)
    
    t("Chapter 9: Learning", "学习与适应模式")
    c("概述", ["从交互中学习", "适应用户偏好", "持续改进"])
    
    prs.save("d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter9_learning/Chapter9_Learning.pptx")
    print("Chapter 9 PPT created")

if __name__ == "__main__":
    create()
