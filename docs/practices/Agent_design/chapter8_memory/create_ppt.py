from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    TC = RGBColor(31, 78, 121)
    
    def t(s, s2=""):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold = s, Pt(44), True
        p.font.color.rgb, p.alignment = TC, PP_ALIGN.CENTER
        if s2:
            sb = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            p = sb.text_frame.paragraphs[0]
            p.text, p.font.size = s2, Pt(24)
            p.alignment = PP_ALIGN.CENTER
    
    def c(t, bs):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text, p.font.size, p.font.bold = t, Pt(32), True
        p.font.color.rgb = TC
        cb = sl.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = cb.text_frame
        for i, b in enumerate(bs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text, p.font.size = "• " + b, Pt(20)
            p.space_after = Pt(12)
    
    t("Chapter 8: Memory Management", "记忆管理 - 让 Agent 记住一切")
    c("什么是记忆管理？", ["短期记忆：当前对话上下文", "长期记忆：持久化存储", "会话间记忆：跨会话信息"])
    c("实现方式", ["List 存储对话历史", "Dict 存储结构化信息", "向量数据库存储语义记忆"])
    c("运行", ["python src/agents/patterns/memory.py"])
    
    prs.save("d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter8_memory/Chapter8_Memory.pptx")
    print("Chapter 8 PPT created")

if __name__ == "__main__":
    create()
