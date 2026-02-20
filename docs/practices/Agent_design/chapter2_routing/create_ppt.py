"""
Chapter 2: Routing 演示文稿生成脚本
使用 python-pptx 库创建 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 颜色定义 - Berry & Cream
colors = {
    'primary': RGBColor(0x6D, 0x2E, 0x46),      # berry
    'secondary': RGBColor(0xA2, 0x67, 0x69),    # dusty rose
    'accent': RGBColor(0xEC, 0xE2, 0xD0),       # cream
    'dark': RGBColor(0x2F, 0x3C, 0x7E),         # navy
    'light': RGBColor(0xFC, 0xF6, 0xF5),        # off-white
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x64, 0x74, 0x8B),
    'text_gray': RGBColor(0x47, 0x55, 0x69),
    'code_bg': RGBColor(0x1E, 0x29, 0x3B),
    'code_text': RGBColor(0xE2, 0xE8, 0xF0),
    'code_highlight': RGBColor(0x22, 0xD3, 0xEE),
    'arrow_gray': RGBColor(0x94, 0xA3, 0xB8),
    'highlight': RGBColor(0x99, 0x00, 0x11),   # cherry
}

def add_text_box(slide, left, top, width, height, text, font_size=12, font_name='Arial', 
                 bold=False, color=None, align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP):
    """添加文本框"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rectangle(slide, left, top, width, height, fill_color):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ============ Slide 1: 封面 ============
slide_layout = prs.slide_layouts[6]  # 空白布局
slide1 = prs.slides.add_slide(slide_layout)

# 背景
background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
background.fill.solid()
background.fill.fore_color.rgb = colors['dark']
background.line.fill.background()

add_text_box(slide1, 0.5, 1.5, 9, 0.8, "Chapter 2", 24, 'Arial', False, colors['accent'])
add_text_box(slide1, 0.5, 2.2, 9, 1.2, "Routing", 48, 'Arial Black', True, colors['white'])
add_text_box(slide1, 0.5, 3.3, 9, 0.6, "路由模式设计模式", 20, 'Arial', False, colors['accent'])

# 装饰线
add_rectangle(slide1, 0.5, 4.2, 2, 0.08, colors['accent'])

add_text_box(slide1, 0.5, 4.5, 9, 0.4, "Everything About Agent - Agent 学习与开发系统", 12, 'Arial', False, colors['arrow_gray'])

# ============ Slide 2: 什么是 Routing ============
slide2 = prs.slides.add_slide(slide_layout)

# 背景
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg2.fill.solid()
bg2.fill.fore_color.rgb = colors['light']
bg2.line.fill.background()

# 顶部装饰条
add_rectangle(slide2, 0, 0, 10, 0.15, colors['primary'])

add_text_box(slide2, 0.5, 0.5, 9, 0.7, "什么是 Routing?", 32, 'Arial Black', True, colors['dark'])

# 装饰线
add_rectangle(slide2, 0.5, 1.1, 1.5, 0.06, colors['highlight'])

# 核心定义卡片
card1 = add_rounded_rectangle(slide2, 0.5, 1.5, 9, 1.3, colors['white'])

definition_text = "Routing（路由模式）引入条件逻辑，使 Agent 能够根据特定条件动态选择不同的处理路径。它将固定的执行流程转变为能够根据输入或状态进行智能决策的灵活系统。"
add_text_box(slide2, 0.7, 1.7, 8.6, 1, definition_text, 16, 'Arial', False, colors['text_gray'])

# 图解标题
add_text_box(slide2, 0.5, 3.1, 9, 0.4, "工作流程示意", 18, 'Arial', True, colors['dark'])

# 路由图 - 中心是 Router，分出三条路径
# 中心 Router
router_box = add_rounded_rectangle(slide2, 4, 3.6, 2, 0.8, colors['primary'])
add_text_box(slide2, 4, 3.6, 2, 0.8, "Router", 14, 'Arial', True, colors['white'], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# 输入
input_box = add_rounded_rectangle(slide2, 0.5, 3.7, 1.5, 0.6, colors['accent'])
add_text_box(slide2, 0.5, 3.7, 1.5, 0.6, "用户输入", 12, 'Arial', True, colors['dark'], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# 箭头到Router
arrow1 = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.1), Inches(3.9), Inches(1.7), Inches(0.25))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = colors['arrow_gray']
arrow1.line.fill.background()

# 三条输出路径
handlers = [
    ("Handler A", "预订处理", colors['secondary']),
    ("Handler B", "信息查询", colors['secondary']),
    ("Handler C", "其他请求", colors['secondary'])
]

for i, (name, desc, color) in enumerate(handlers):
    x = 7 + i * 0.95
    y = 3.2 + i * 0.4
    
    # 小方块
    add_rounded_rectangle(slide2, x, y, 1.8, 0.5, color)
    add_text_box(slide2, x, y, 1.8, 0.5, name, 10, 'Arial', True, colors['white'], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ============ Slide 3: 核心实现方式 ============
slide3 = prs.slides.add_slide(slide_layout)

bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg3.fill.solid()
bg3.fill.fore_color.rgb = colors['light']
bg3.line.fill.background()

add_rectangle(slide3, 0, 0, 10, 0.15, colors['primary'])
add_text_box(slide3, 0.5, 0.5, 9, 0.7, "核心实现方式", 32, 'Arial Black', True, colors['dark'])
add_rectangle(slide3, 0.5, 1.1, 1.5, 0.06, colors['highlight'])

# 四种实现方式
methods = [
    ("LLM-based Routing", "使用 LLM 分析输入并输出分类标识", "🤖"),
    ("Embedding-based Routing", "将输入转换为向量，通过语义相似度匹配路由", "📊"),
    ("Rule-based Routing", "使用预定义规则（关键词、模式）进行路由", "📏"),
    ("ML Model-based Routing", "使用训练好的分类模型进行路由决策", "🧠")
]

for i, (title, desc, icon) in enumerate(methods):
    x = 0.5 + (i % 2) * 4.8
    y = 1.5 + (i // 2) * 1.8
    
    # 卡片背景
    add_rounded_rectangle(slide3, x, y, 4.5, 1.5, colors['white'])
    
    # 左侧装饰条
    add_rectangle(slide3, x, y, 0.1, 1.5, colors['primary'])
    
    # 图标
    add_text_box(slide3, x + 0.2, y + 0.15, 0.6, 0.5, icon, 24, 'Arial', False, None, PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide3, x + 0.8, y + 0.2, 3.5, 0.4, title, 14, 'Arial', True, colors['dark'])
    
    # 描述
    desc_box = slide3.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.65), Inches(3.5), Inches(0.7))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.name = 'Arial'
    p.font.color.rgb = colors['gray']

# ============ Slide 4: 实际应用场景 ============
slide4 = prs.slides.add_slide(slide_layout)

bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg4.fill.solid()
bg4.fill.fore_color.rgb = colors['light']
bg4.line.fill.background()

add_rectangle(slide4, 0, 0, 10, 0.15, colors['secondary'])
add_text_box(slide4, 0.5, 0.5, 9, 0.7, "实际应用场景", 28, 'Arial Black', True, colors['dark'])

# 场景卡片
scenarios = [
    ("客户服务机器人", "区分订单查询、产品咨询、技术支持等不同意图", "🛎️"),
    ("数据处理管道", "根据数据类型分发到不同处理流程", "📧"),
    ("多工具调度", "根据任务需求选择最合适的工具或子Agent", "🔧"),
    ("编程助手", "识别编程语言和意图（调试、解释、翻译）", "💻")
]

for i, (title, desc, icon) in enumerate(scenarios):
    x = 0.5 + (i % 2) * 4.8
    y = 1.4 + (i // 2) * 1.9
    
    # 卡片
    add_rounded_rectangle(slide4, x, y, 4.5, 1.6, colors['white'])
    
    # 顶部装饰
    add_rectangle(slide4, x, y, 4.5, 0.08, colors['primary'])
    
    # 图标和标题
    add_text_box(slide4, x + 0.15, y + 0.2, 4.2, 0.5, f"{icon}  {title}", 14, 'Arial', True, colors['dark'])
    
    # 描述
    desc_box = slide4.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.75), Inches(4.2), Inches(0.7))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.name = 'Arial'
    p.font.color.rgb = colors['text_gray']

# ============ Slide 5: LangChain 实现 ============
slide5 = prs.slides.add_slide(slide_layout)

bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg5.fill.solid()
bg5.fill.fore_color.rgb = colors['light']
bg5.line.fill.background()

add_rectangle(slide5, 0, 0, 10, 0.15, colors['secondary'])
add_text_box(slide5, 0.5, 0.5, 9, 0.7, "LangChain 实现示例", 28, 'Arial Black', True, colors['dark'])
add_text_box(slide5, 0.5, 1.05, 9, 0.4, "使用 RunnableBranch 实现路由", 14, 'Arial', False, colors['gray'])

# 代码示例
code_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = colors['code_bg']
code_bg.line.fill.background()

code_text = '''# 定义路由提示
router_prompt = ChatPromptTemplate.from_messages([
    ("system", "分析用户请求，输出分类: 'booker', 'info', 或 'unclear'"),
    ("user", "{request}")
])

# 创建路由链
router_chain = router_prompt | llm | StrOutputParser()

# 使用 RunnableBranch 实现分支逻辑
branch = RunnableBranch(
    (lambda x: "booker" in x["decision"], booking_handler),
    (lambda x: "info" in x["decision"], info_handler),
    unclear_handler  # 默认分支
)

# 组合完整链
full_chain = {"decision": router_chain, "request": RunnablePassthrough()} | branch'''

code_box = slide5.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(3.3))
tf = code_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(10)
p.font.name = 'Consolas'
p.font.color.rgb = colors['code_text']

# ============ Slide 6: 项目代码示例 ============
slide6 = prs.slides.add_slide(slide_layout)

bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg6.fill.solid()
bg6.fill.fore_color.rgb = colors['light']
bg6.line.fill.background()

add_rectangle(slide6, 0, 0, 10, 0.15, colors['highlight'])
add_text_box(slide6, 0.5, 0.5, 9, 0.7, "项目代码示例", 28, 'Arial Black', True, colors['dark'])
add_text_box(slide6, 0.5, 1.05, 9, 0.4, "src/agents/patterns/routing.py", 14, 'Arial', False, colors['gray'])

# 场景说明
add_rounded_rectangle(slide6, 0.5, 1.5, 4.3, 1.2, colors['white'])
add_text_box(slide6, 0.7, 1.6, 3.9, 0.35, "场景", 14, 'Arial', True, colors['primary'])
add_text_box(slide6, 0.7, 2.0, 3.9, 0.6, "根据用户意图，将请求路由到不同的处理程序（预订 vs 信息查询）", 12, 'Arial', False, colors['text_gray'])

# 流程标题
add_text_box(slide6, 5.2, 1.5, 4.3, 0.4, "处理流程", 14, 'Arial', True, colors['dark'])

# 流程步骤
steps = [
    ("输入", "用户请求", colors['accent']),
    ("Router", "LLM 分析意图", colors['primary']),
    ("分支", "RunnableBranch", colors['secondary']),
    ("输出", "Handler 结果", colors['accent'])
]

for i, (title, desc, color) in enumerate(steps):
    y = 2.0 + i * 0.55
    
    add_rounded_rectangle(slide6, 5.2, y, 1.1, 0.45, color)
    add_text_box(slide6, 5.2, y, 1.1, 0.45, title, 10, 'Arial', True, colors['dark'] if color == colors['accent'] else colors['white'], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text_box(slide6, 6.4, y, 3.1, 0.45, desc, 11, 'Arial', False, colors['text_gray'])

# 运行说明
add_text_box(slide6, 0.5, 3.0, 4.3, 0.4, "如何运行", 14, 'Arial', True, colors['dark'])

# 代码背景
code_bg2 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.4), Inches(9), Inches(1.4))
code_bg2.fill.solid()
code_bg2.fill.fore_color.rgb = colors['code_bg']
code_bg2.line.fill.background()

add_text_box(slide6, 0.6, 3.5, 8.8, 1.2, "python src/agents/patterns/routing.py\n# 或在 start.bat 中选择选项 4", 11, 'Consolas', False, colors['code_text'])

# ============ Slide 7: 总结 ============
slide7 = prs.slides.add_slide(slide_layout)

bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg7.fill.solid()
bg7.fill.fore_color.rgb = colors['dark']
bg7.line.fill.background()

add_text_box(slide7, 0.5, 0.8, 9, 0.7, "总结", 36, 'Arial Black', True, colors['white'])
add_rectangle(slide7, 0.5, 1.4, 2, 0.06, colors['accent'])

# 关键要点
key_points = [
    "Routing 使 Agent 能够根据条件动态选择执行路径",
    "支持 LLM-based、Embedding-based、Rule-based 等多种实现方式",
    "LangChain 使用 RunnableBranch 实现分支逻辑",
    "适合需要多任务分发的智能系统场景"
]

for i, point in enumerate(key_points):
    # 圆点
    dot = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.8 + i * 0.7), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors['accent']
    dot.line.fill.background()
    
    add_text_box(slide7, 1.0, 1.75 + i * 0.7, 8.5, 0.5, point, 14, 'Arial', False, colors['white'])

# 底部
add_text_box(slide7, 0.5, 4.5, 9, 0.4, "下一章: Parallelization (并行模式)", 14, 'Arial', False, colors['accent'])

# 保存文件
output_path = r"d:\code\python\everything_about_agent\docs\practices\Agent_design\chapter2_routing\Chapter2_Routing.pptx"
prs.save(output_path)
print(f"PPT 已保存到: {output_path}")
