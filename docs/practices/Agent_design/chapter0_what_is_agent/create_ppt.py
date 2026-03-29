"""
创建"什么是AI Agent"的PPT演示文稿
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案 - 科技蓝主题
COLORS = {
    'primary_dark': RGBColor(0x15, 0x1B, 0x2B),      # 深蓝背景
    'primary': RGBColor(0x3B, 0x82, 0xF6),            # 亮蓝
    'secondary': RGBColor(0x60, 0xA5, 0xFA),          # 浅蓝
    'accent': RGBColor(0xF5, 0x95, 0x0B),             # 橙色强调
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xE5, 0xE7, 0xEB),
    'text_gray': RGBColor(0x9C, 0xA3, 0xAF),
    'green': RGBColor(0x10, 0xB9, 0x81),              # 绿色
    'purple': RGBColor(0x8B, 0x5C, 0xF6),             # 紫色
}

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary_dark']
    background.line.fill.background()
    
    # 装饰线条
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Microsoft YaHei"
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['secondary']
        p.font.name = "Microsoft YaHei"
    
    return slide

def add_content_slide(prs, title, bullets, note=""):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 顶部装饰条
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # 左侧装饰条
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), prs.slide_height)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['primary_dark']
    left_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Microsoft YaHei"
    
    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.level = 0
    
    # 演讲者备注
    if note:
        slide.notes_slide.notes_text_frame.text = note
    
    return slide

def add_comparison_slide(prs, title, headers, rows):
    """添加对比表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 装饰
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Microsoft YaHei"
    
    # 创建表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5)).table
    
    # 设置列宽 (需要整数EMU单位)
    col_width = int(Inches(12).emu / num_cols)
    for col in table.columns:
        col.width = col_width
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = COLORS['white']
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Microsoft YaHei"
        paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 数据行
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Microsoft YaHei"
            paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 隔行变色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
    
    return slide

def add_architecture_slide(prs):
    """添加架构图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary_dark']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Agent 架构组成"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    
    # 三个主要模块
    modules = [
        ("感知模块\nPerception", Inches(1), COLORS['secondary'], 
         ["收集环境信息", "自然语言理解", "多模态输入"]),
        ("大脑/控制端\nBrain / LLM", Inches(5.2), COLORS['primary'],
         ["推理与决策", "任务规划", "反思与改进"]),
        ("行动模块\nAction", Inches(9.4), COLORS['green'],
         ["调用工具", "执行操作", "环境交互"]),
    ]
    
    for title, left, color, desc_list in modules:
        # 模块框
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(3), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # 模块标题
        title_tb = slide.shapes.add_textbox(left, Inches(2.2), Inches(3), Inches(0.8))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        
        # 模块描述
        desc_tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(3), Inches(2.6), Inches(2.3))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        for i, desc in enumerate(desc_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {desc}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['white']
            p.font.name = "Microsoft YaHei"
            p.space_before = Pt(8)
    
    # 连接箭头（用简单线条表示）
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(1.2), Pt(4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLORS['light_gray']
    arrow1.line.fill.background()
    
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(3.5), Inches(1.2), Pt(4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS['light_gray']
    arrow2.line.fill.background()
    
    # 记忆模块（在下方）
    memory_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(5.8), Inches(3), Inches(1.2))
    memory_box.fill.solid()
    memory_box.fill.fore_color.rgb = COLORS['purple']
    memory_box.line.fill.background()
    
    memory_tb = slide.shapes.add_textbox(Inches(5.2), Inches(6), Inches(3), Inches(0.8))
    tf = memory_tb.text_frame
    p = tf.paragraphs[0]
    p.text = "记忆模块 Memory\n短期记忆 + 长期记忆"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    
    # 连接记忆
    mem_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(5.5), Pt(4), Inches(0.3))
    mem_line.fill.solid()
    mem_line.fill.fore_color.rgb = COLORS['purple']
    mem_line.line.fill.background()
    
    return slide

def add_workflow_slide(prs):
    """添加工作流程图页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 浅色背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Agent 工作循环：感知-思考-行动"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    
    # 四个步骤
    steps = [
        ("感知\nPerceive", Inches(0.8), COLORS['secondary'], "👁️"),
        ("思考\nThink", Inches(3.8), COLORS['primary'], "🧠"),
        ("行动\nAct", Inches(6.8), COLORS['green'], "🛠️"),
        ("反馈\nFeedback", Inches(9.8), RGBColor(0xEC, 0x48, 0x99), "🔄"),
    ]
    
    for i, (title, left, color, emoji) in enumerate(steps):
        # 圆形
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.5), Inches(1.5), Inches(1.5), Inches(1.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        # 步骤标题
        title_tb = slide.shapes.add_textbox(left, Inches(3.1), Inches(2.5), Inches(0.8))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        
        # 箭头（除了最后一个）
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(2), Inches(2.1), Inches(1.3), Pt(6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['text_gray']
            arrow.line.fill.background()
    
    # 详细说明
    details = [
        (Inches(0.5), "• 接收用户指令\n• 收集环境信息\n• 多模态输入处理", COLORS['secondary']),
        (Inches(3.5), "• 分析问题\n• 制定计划\n• 分解任务", COLORS['primary']),
        (Inches(6.5), "• 调用工具\n• 执行操作\n• 与环境交互", COLORS['green']),
        (Inches(9.5), "• 观察结果\n• 学习改进\n• 循环优化", RGBColor(0xEC, 0x48, 0x99)),
    ]
    
    for left, text, color in details:
        detail_tb = slide.shapes.add_textbox(left, Inches(4), Inches(2.8), Inches(2.5))
        tf = detail_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
    
    return slide

def add_analogy_slide(prs):
    """添加类比说明页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 装饰条
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Agent vs LLM vs RAG：一个形象的比喻"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Microsoft YaHei"
    
    # 三个对比卡片
    cards = [
        ("LLM\n大语言模型", Inches(0.6), RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0xDC, 0x26, 0x26),
         "🧠", "聪明的大脑被困在\n玻璃罐里", "• 知识渊博\n• 无法感知外部世界\n• 不能执行操作"),
        ("RAG\n检索增强", Inches(4.6), RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xD9, 0x77, 0x06),
         "📚", "大脑+图书管理员", "• 能获取外部知识\n• 信息更准确\n• 仍不能采取行动"),
        ("AI Agent\n智能代理", Inches(8.6), RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0x05, 0x96, 0x69),
         "🤖", "大脑有了完整身体", "• 能感知世界\n• 能使用工具\n• 能自主决策行动"),
    ]
    
    for title, left, bg_color, title_color, emoji, subtitle, features in cards:
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.fill.background()
        
        # 标题区域
        title_area = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(0.8))
        title_area.fill.solid()
        title_area.fill.fore_color.rgb = title_color
        title_area.line.fill.background()
        
        # 标题文字
        title_tb = slide.shapes.add_textbox(left, Inches(1.6), Inches(3.8), Inches(0.6))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_tb = slide.shapes.add_textbox(left, Inches(2.5), Inches(3.8), Inches(0.6))
        tf = sub_tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        
        # 特性列表
        feat_tb = slide.shapes.add_textbox(left + Inches(0.3), Inches(3.2), Inches(3.2), Inches(3))
        tf = feat_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = features
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.font.name = "Microsoft YaHei"
        p.line_spacing = 1.5
    
    return slide

# ==================== 开始创建幻灯片 ====================

# 1. 封面
add_title_slide(prs, "什么是 AI Agent？", "人工智能代理基础概念与架构解析")

# 2. 学习目标
add_content_slide(prs, "🎯 本章学习目标", [
    "理解 AI Agent 的核心定义与本质特征",
    "掌握 AI Agent 与 LLM、RAG 的区别与联系",
    "了解 AI Agent 的核心架构组成",
    "理解 AI Agent 的工作循环与决策流程",
    "认识 AI Agent 的典型应用场景"
])

# 3. 什么是 AI Agent
add_content_slide(prs, "1. 什么是 AI Agent？", [
    "AI Agent（人工智能代理）是能够自主感知环境、制定决策并执行行动的智能系统",
    "核心特点：自主性（Autonomous）、主动性（Proactive）、目标导向（Goal-oriented）",
    "与传统 AI 的区别：不只是问答，而是能够主动思考、规划并完成复杂任务",
    "权威定义：Agent = LLM + 记忆能力 + 规划能力 + 工具使用能力（OpenAI）",
    '简单类比：如果传统 AI 是"问答机器"，AI Agent 就是"智能助手/数字员工"'
])

# 4. 类比说明
add_analogy_slide(prs)

# 5. 核心特征
add_content_slide(prs, "2. AI Agent 的核心特征", [
    "自主性（Autonomy）：能够独立做出决策，无需人类持续干预",
    "反应性（Reactivity）：能够感知环境变化并及时响应",
    "主动性（Proactivity）：能够主动采取行动实现目标",
    "社交性（Social Ability）：能够与其他 Agent 或人类进行交互协作",
    "与传统 RPA 的区别：AI Agent 能够动态调整、从经验中学习改进"
])

# 6. 架构组成
add_architecture_slide(prs)

# 7. 大脑模块详解
add_content_slide(prs, "3.1 大脑模块（Brain / LLM）", [
    "核心组件：大语言模型（GPT-4、Claude、Llama、Qwen 等）",
    "推理与决策：分析问题、制定策略、选择最优方案",
    "规划能力：将复杂任务分解为可执行的子任务",
    "反思与改进：从错误中学习，优化后续行动",
    "自然语言交互：理解指令、生成回复、与用户沟通"
])

# 8. 感知与行动模块
add_content_slide(prs, "3.2 感知模块 & 行动模块", [
    "感知模块（Perception）：",
    "   • 收集和处理环境信息：文本、图像、音频、结构化数据",
    "   • 自然语言理解、计算机视觉、传感器数据处理",
    "",
    "行动模块（Action）：",
    "   • 执行决策、与环境交互",
    "   • 工具使用：搜索引擎、API调用、代码执行、文件操作",
    "   • 输出：文本回复、操作执行、状态更新"
])

# 9. 记忆模块
add_content_slide(prs, "3.3 记忆模块（Memory）", [
    "短期记忆（Short-term Memory）：",
    "   • 当前对话上下文、任务执行状态",
    "   • 临时工作记忆，会话结束后清除",
    "",
    "长期记忆（Long-term Memory）：",
    "   • 用户偏好与历史交互记录",
    "   • 知识库与经验积累",
    "   • 通常使用向量数据库存储，支持快速检索"
])

# 10. 工作循环
add_workflow_slide(prs)

# 11. ReAct 模式
add_content_slide(prs, "4. ReAct 模式：推理与行动的融合", [
    "ReAct = Reasoning（推理）+ Acting（行动）",
    "工作循环：Thought → Action → Observation → Thought → ...",
    '让 AI Agent 像人类一样"一步步思考"，提高复杂问题解决能力',
    "示例：查询天气 → 搜索信息 → 分析结果 → 给出建议",
    "核心优势：能够迭代推理、自我纠错、动态调整策略"
])

# 12. 对比表格
add_comparison_slide(prs, "5. AI Agent vs LLM vs RAG 对比", 
    ["能力", "LLM", "RAG", "AI Agent"],
    [
        ["文本生成", "✅", "✅", "✅"],
        ["访问外部文档", "❌", "✅", "✅"],
        ["执行代码", "❌", "❌", "✅"],
        ["调用 API", "❌", "❌", "✅"],
        ["迭代推理", "有限", "❌", "✅"],
        ["自我纠错", "❌", "❌", "✅"],
        ["多步规划", "❌", "❌", "✅"],
        ["自主运行", "❌", "❌", "✅"],
    ]
)

# 13. 典型应用场景
add_content_slide(prs, "6. AI Agent 的典型应用场景", [
    "智能客服 Agent：理解问题、查询知识库、执行操作、转接人工",
    "编程助手 Agent（如 Cursor、Devin）：编写、调试、测试代码",
    "研究助手 Agent：自主搜索、分析论文、生成报告",
    "个人助理 Agent：日程管理、邮件处理、旅行规划",
    "数据分析 Agent：自动处理数据、生成可视化、提供业务洞察"
])

# 14. 技术栈
add_content_slide(prs, "7. AI Agent 技术栈", [
    "模型层：GPT-4、Claude、Llama、Qwen 等大语言模型",
    "开发框架：LangChain（模块化）、AutoGen（多Agent）、CrewAI（角色扮演）",
    "记忆存储：向量数据库（Chroma、Pinecone）、知识图谱",
    "工具集成：搜索引擎、API网关、代码执行环境、浏览器自动化",
    "基础设施：云服务器、容器化部署、监控与日志系统"
])

# 15. 发展趋势
add_content_slide(prs, "8. 发展趋势与展望", [
    "从单 Agent 到多 Agent 协作：多个专业 Agent 组成团队",
    "从工具使用到深度集成：MCP（Model Context Protocol）标准化",
    "从通用到垂直领域：法律、医疗、金融等专业 Agent",
    "从云端到端侧部署：轻量级模型支持本地运行",
    "未来愿景：通用数字员工 → 物理世界机器人 → AGI 雏形"
])

# 16. 总结
add_content_slide(prs, "📌 本章小结", [
    "AI Agent 本质：LLM + 记忆 + 规划 + 工具使用能力",
    "核心特征：自主性、反应性、主动性、社交性",
    "架构组成：感知 → 大脑（LLM）→ 行动 + 记忆",
    "工作循环：感知 → 思考 → 行动 → 反馈",
    "关键区别：Agent 能够 思考并行动，而不仅是回答问题"
])

# 17. 结尾页
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 深色背景
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = COLORS['primary_dark']
background.line.fill.background()

# 标题
title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "感谢观看"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

# 副标题
sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "下一章：Prompt Chaining（提示链）"
p.font.size = Pt(24)
p.font.color.rgb = COLORS['secondary']
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

# 保存
output_path = "docs/practices/Agent_design/chapter0_what_is_agent/Chapter0_什么是AIAgent.pptx"
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
