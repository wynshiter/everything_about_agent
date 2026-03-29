"""
Create "What is AI Agent" PowerPoint Presentation (English Version)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Tech Blue theme
COLORS = {
    'primary_dark': RGBColor(0x15, 0x1B, 0x2B),
    'primary': RGBColor(0x3B, 0x82, 0xF6),
    'secondary': RGBColor(0x60, 0xA5, 0xFA),
    'accent': RGBColor(0xF5, 0x95, 0x0B),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xE5, 0xE7, 0xEB),
    'text_gray': RGBColor(0x9C, 0xA3, 0xAF),
    'green': RGBColor(0x10, 0xB9, 0x81),
    'purple': RGBColor(0x8B, 0x5C, 0xF6),
}

def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary_dark']
    background.line.fill.background()
    
    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['primary']
    line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Arial"
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['secondary']
        p.font.name = "Arial"
    
    return slide

def add_content_slide(prs, title, bullets, note=""):
    """Add content slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # Top decoration
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # Left decoration
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), prs.slide_height)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['primary_dark']
    left_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Arial"
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.font.name = "Arial"
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.level = 0
    
    if note:
        slide.notes_slide.notes_text_frame.text = note
    
    return slide

def add_comparison_slide(prs, title, headers, rows):
    """Add comparison table slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # Decoration
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Arial"
    
    # Create table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5)).table
    
    # Set column width
    col_width = int(Inches(12).emu / num_cols)
    for col in table.columns:
        col.width = col_width
    
    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = COLORS['white']
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Arial"
            paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)
    
    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary_dark']
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Agent Architecture"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # Three main modules
    modules = [
        ("Perception\nModule", Inches(1), COLORS['secondary'], 
         ["Collect environment info", "Natural language understanding", "Multi-modal input"]),
        ("Brain / Control\n(LLM)", Inches(5.2), COLORS['primary'],
         ["Reasoning & decision", "Task planning", "Reflection & improvement"]),
        ("Action\nModule", Inches(9.4), COLORS['green'],
         ["Call tools", "Execute operations", "Environment interaction"]),
    ]
    
    for title, left, color, desc_list in modules:
        # Module box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(3), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Module title
        title_tb = slide.shapes.add_textbox(left, Inches(2.2), Inches(3), Inches(0.8))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        # Module description
        desc_tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(3), Inches(2.6), Inches(2.3))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        for i, desc in enumerate(desc_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {desc}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['white']
            p.font.name = "Arial"
            p.space_before = Pt(8)
    
    # Connection arrows
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(1.2), Pt(4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLORS['light_gray']
    arrow1.line.fill.background()
    
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(3.5), Inches(1.2), Pt(4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS['light_gray']
    arrow2.line.fill.background()
    
    # Memory module (below)
    memory_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(5.8), Inches(3), Inches(1.2))
    memory_box.fill.solid()
    memory_box.fill.fore_color.rgb = COLORS['purple']
    memory_box.line.fill.background()
    
    memory_tb = slide.shapes.add_textbox(Inches(5.2), Inches(6), Inches(3), Inches(0.8))
    tf = memory_tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Memory Module\nShort-term + Long-term"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # Connect memory
    mem_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(5.5), Pt(4), Inches(0.3))
    mem_line.fill.solid()
    mem_line.fill.fore_color.rgb = COLORS['purple']
    mem_line.line.fill.background()
    
    return slide

def add_workflow_slide(prs):
    """Add workflow diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Agent Working Loop: Perceive-Think-Act"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    # Four steps
    steps = [
        ("Perceive", Inches(0.8), COLORS['secondary'], "👁️"),
        ("Think", Inches(3.8), COLORS['primary'], "🧠"),
        ("Act", Inches(6.8), COLORS['green'], "🛠️"),
        ("Feedback", Inches(9.8), RGBColor(0xEC, 0x48, 0x99), "🔄"),
    ]
    
    for i, (title, left, color, emoji) in enumerate(steps):
        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.5), Inches(1.5), Inches(1.5), Inches(1.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        # Step title
        title_tb = slide.shapes.add_textbox(left, Inches(3.1), Inches(2.5), Inches(0.8))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(2), Inches(2.1), Inches(1.3), Pt(6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['text_gray']
            arrow.line.fill.background()
    
    # Detailed description
    details = [
        (Inches(0.5), "• Receive user input\n• Collect environment info\n• Multi-modal processing", COLORS['secondary']),
        (Inches(3.5), "• Analyze problem\n• Make plan\n• Decompose tasks", COLORS['primary']),
        (Inches(6.5), "• Call tools\n• Execute operations\n• Interact with environment", COLORS['green']),
        (Inches(9.5), "• Observe results\n• Learn & improve\n• Iterate optimization", RGBColor(0xEC, 0x48, 0x99)),
    ]
    
    for left, text, color in details:
        detail_tb = slide.shapes.add_textbox(left, Inches(4), Inches(2.8), Inches(2.5))
        tf = detail_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = color
        p.font.name = "Arial"
    
    return slide

def add_analogy_slide(prs):
    """Add analogy slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # Decoration
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "LLM vs RAG vs AI Agent: An Analogy"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_dark']
    p.font.name = "Arial"
    
    # Three comparison cards
    cards = [
        ("LLM\nLanguage Model", Inches(0.6), RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0xDC, 0x26, 0x26),
         "🧠", "Brilliant brain in\na sealed jar", "• Vast knowledge\n• Cannot perceive world\n• Cannot take actions"),
        ("RAG\nRetrieval Augmented", Inches(4.6), RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xD9, 0x77, 0x06),
         "📚", "Brain + Librarian", "• Access external knowledge\n• More accurate info\n• Still cannot act"),
        ("AI Agent\nAutonomous Agent", Inches(8.6), RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0x05, 0x96, 0x69),
         "🤖", "Brain with a\nfull body", "• Can perceive world\n• Can use tools\n• Can act autonomously"),
    ]
    
    for title, left, bg_color, title_color, emoji, subtitle, features in cards:
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.fill.background()
        
        # Title area
        title_area = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(0.8))
        title_area.fill.solid()
        title_area.fill.fore_color.rgb = title_color
        title_area.line.fill.background()
        
        # Title text
        title_tb = slide.shapes.add_textbox(left, Inches(1.6), Inches(3.8), Inches(0.6))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_tb = slide.shapes.add_textbox(left, Inches(2.5), Inches(3.8), Inches(0.6))
        tf = sub_tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        # Feature list
        feat_tb = slide.shapes.add_textbox(left + Inches(0.3), Inches(3.2), Inches(3.2), Inches(3))
        tf = feat_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = features
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.font.name = "Arial"
        p.line_spacing = 1.5
    
    return slide

# ==================== Start Creating Slides ====================

# 1. Cover
add_title_slide(prs, "What is AI Agent?", "Fundamental Concepts and Architecture")

# 2. Learning Objectives
add_content_slide(prs, "🎯 Learning Objectives", [
    "Understand the core definition and essential characteristics of AI Agent",
    "Master the differences and connections between AI Agent, LLM, and RAG",
    "Learn the core architecture components of AI Agent",
    "Understand the working loop and decision-making process of AI Agent",
    "Recognize typical application scenarios of AI Agent"
])

# 3. What is AI Agent
add_content_slide(prs, "1. What is AI Agent?", [
    "AI Agent is an intelligent system that can perceive environment, make decisions, and take actions autonomously",
    "Core characteristics: Autonomous, Proactive, Goal-oriented",
    "Difference from traditional AI: Not just Q&A, but active thinking, planning, and completing complex tasks",
    "Authoritative definition: Agent = LLM + Memory + Planning + Tool Use (OpenAI)",
    'Simple analogy: If traditional AI is a "Q&A machine", AI Agent is an "intelligent assistant/digital employee"'
])

# 4. Analogy
add_analogy_slide(prs)

# 5. Core Characteristics
add_content_slide(prs, "2. Core Characteristics of AI Agent", [
    "Autonomy: Can make decisions independently without continuous human intervention",
    "Reactivity: Can perceive environmental changes and respond promptly",
    "Proactivity: Can take initiative to achieve goals",
    "Social Ability: Can interact and collaborate with other agents or humans",
    "Difference from RPA: AI Agent can dynamically adjust and learn from experience"
])

# 6. Architecture
add_architecture_slide(prs)

# 7. Brain Module
add_content_slide(prs, "3.1 Brain Module (LLM)", [
    "Core component: Large Language Models (GPT-4, Claude, Llama, Qwen, etc.)",
    "Reasoning & Decision: Analyze problems, formulate strategies, select optimal solutions",
    "Planning: Break down complex tasks into executable subtasks",
    "Reflection & Improvement: Learn from mistakes and optimize future actions",
    "Natural Language Interaction: Understand instructions, generate responses, communicate with users"
])

# 8. Perception & Action Module
add_content_slide(prs, "3.2 Perception & Action Module", [
    "Perception Module:",
    "   • Collect and process environment information: text, image, audio, structured data",
    "   • Natural language understanding, computer vision, sensor data processing",
    "",
    "Action Module:",
    "   • Execute decisions, interact with environment",
    "   • Tool use: search engine, API calls, code execution, file operations",
    "   • Output: text response, operation execution, status update"
])

# 9. Memory Module
add_content_slide(prs, "3.3 Memory Module", [
    "Short-term Memory:",
    "   • Current conversation context, task execution status",
    "   • Temporary working memory, cleared after session ends",
    "",
    "Long-term Memory:",
    "   • User preferences and historical interaction records",
    "   • Knowledge base and experience accumulation",
    "   • Usually stored in vector databases for fast retrieval"
])

# 10. Working Loop
add_workflow_slide(prs)

# 11. ReAct Pattern
add_content_slide(prs, "4. ReAct Pattern: Reasoning + Acting", [
    "ReAct = Reasoning + Acting",
    "Working loop: Thought → Action → Observation → Thought → ...",
    'Enable AI Agent to "think step by step" like humans, improving complex problem-solving',
    "Example: Query weather → Search info → Analyze results → Give advice",
    "Core advantages: Iterative reasoning, self-correction, dynamic strategy adjustment"
])

# 12. Comparison Table
add_comparison_slide(prs, "5. AI Agent vs LLM vs RAG Comparison", 
    ["Capability", "LLM", "RAG", "AI Agent"],
    [
        ["Text Generation", "✅", "✅", "✅"],
        ["Access External Docs", "❌", "✅", "✅"],
        ["Execute Code", "❌", "❌", "✅"],
        ["Call API", "❌", "❌", "✅"],
        ["Iterative Reasoning", "Limited", "❌", "✅"],
        ["Self-correction", "❌", "❌", "✅"],
        ["Multi-step Planning", "❌", "❌", "✅"],
        ["Autonomous Operation", "❌", "❌", "✅"],
    ]
)

# 13. Application Scenarios
add_content_slide(prs, "6. Typical Application Scenarios", [
    "Intelligent Customer Service: Understand questions, query knowledge base, execute operations",
    "Coding Assistant (Cursor, Devin): Write, debug, test code autonomously",
    "Research Assistant: Autonomous search, analyze papers, generate reports",
    "Personal Assistant: Schedule management, email processing, travel planning",
    "Data Analysis Agent: Automatically process data, generate visualizations, provide insights"
])

# 14. Tech Stack
add_content_slide(prs, "7. AI Agent Technology Stack", [
    "Model Layer: GPT-4, Claude, Llama, Qwen and other LLMs",
    "Frameworks: LangChain (modular), AutoGen (multi-agent), CrewAI (role-playing)",
    "Memory Storage: Vector databases (Chroma, Pinecone), Knowledge graphs",
    "Tool Integration: Search engines, API gateways, code execution, browser automation",
    "Infrastructure: Cloud servers, container deployment, monitoring and logging"
])

# 15. Development Trends
add_content_slide(prs, "8. Trends and Outlook", [
    "From Single Agent to Multi-Agent Collaboration: Multiple specialized agents form teams",
    "From Tool Use to Deep Integration: MCP (Model Context Protocol) standardization",
    "From General to Vertical Domains: Legal, medical, financial specialized agents",
    "From Cloud to Edge Deployment: Lightweight models support local execution",
    "Future Vision: General Digital Employee → Physical World Robot → AGI Prototype"
])

# 16. Summary
add_content_slide(prs, "📌 Summary", [
    "AI Agent essence: LLM + Memory + Planning + Tool Use",
    "Core characteristics: Autonomy, Reactivity, Proactivity, Social Ability",
    "Architecture: Perception → Brain (LLM) → Action + Memory",
    "Working loop: Perceive → Think → Act → Feedback",
    "Key difference: Agent can think AND act, not just answer questions"
])

# 17. End slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Dark background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = COLORS['primary_dark']
background.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.font.name = "Arial"
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Next: Prompt Chaining (提示链)"
p.font.size = Pt(24)
p.font.color.rgb = COLORS['secondary']
p.font.name = "Arial"
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "docs/practices/Agent_design/chapter0_what_is_agent/Chapter0_What_is_AI_Agent_EN.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
