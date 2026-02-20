"""
Chapter 3: Parallelization PPT Generator
生成 Parallelization 模式的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_parallelization_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 颜色定义
    TITLE_COLOR = RGBColor(31, 78, 121)  # 深蓝色
    ACCENT_COLOR = RGBColor(0, 120, 212)  # 蓝色
    
    def add_title_slide(title: str, subtitle: str = ""):
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title: str, bullets: list):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.space_after = Pt(12)
        
        return slide
    
    # Slide 1: 封面
    add_title_slide(
        "Chapter 3: Parallelization",
        "并行化模式 - 提升 Agent 效率的核心技术"
    )
    
    # Slide 2: 什么是并行化
    add_content_slide(
        "什么是并行化模式？",
        [
            "同时执行多个独立的子任务，而不是顺序执行",
            "减少总体执行时间，特别是涉及外部 I/O 操作时",
            "适用于 API 调用、数据库查询、数据处理等场景",
            "核心思想：识别工作流中不依赖其他部分输出的组件并并行执行"
        ]
    )
    
    # Slide 3: 适用场景
    add_content_slide(
        "并行化的典型应用场景",
        [
            "信息收集与研究：同时从多个来源搜索信息",
            "数据处理与分析：对不同数据段并行执行分析",
            "多 API 调用：并行查询天气、新闻、股票等",
            "内容生成：并行生成标题、正文、图片描述",
            "验证与检查：并行验证邮箱、电话、地址格式"
        ]
    )
    
    # Slide 4: LangChain 实现
    add_content_slide(
        "LangChain LCEL 实现",
        [
            "使用 RunnableParallel 并行执行多个 Runnable",
            "使用 RunnablePassthrough 传递原始输入",
            "示例结构：{summary: chain1, questions: chain2, terms: chain3}",
            "分支结果自动聚合后传入下一步进行合成"
        ]
    )
    
    # Slide 5: Google ADK 实现
    add_content_slide(
        "Google ADK 实现",
        [
            "使用 ParallelAgent 定义并行执行的子 Agent",
            "每个子 Agent 可设置 output_key 存储结果到状态",
            "子 Agent 完成后，结果自动聚合到 session state",
            "可配合 SequentialAgent 实现并行+顺序的混合工作流"
        ]
    )
    
    # Slide 6: 代码示例
    add_content_slide(
        "代码示例 - LangChain",
        [
            "from langchain_core.runnables import RunnableParallel",
            "",
            "map_chain = RunnableParallel({",
            '    "summary": summarize_chain,',
            '    "questions": questions_chain,',
            '    "key_terms": terms_chain',
            "})",
            "",
            "full_chain = map_chain | synthesis_prompt | llm | parser"
        ]
    )
    
    # Slide 7: 运行指南
    add_content_slide(
        "运行指南",
        [
            "确保已安装依赖：pip install -e .",
            "确保 Ollama 或其他后端服务正在运行",
            "运行示例：python src/agents/patterns/parallelization.py",
            "查看项目 README 了解更多配置细节"
        ]
    )
    
    # Slide 8: 总结
    add_content_slide(
        "总结与要点",
        [
            "并行化是优化 Agent 性能的关键技术",
            "适用于有外部延迟的独立任务",
            "LangChain 使用 RunnableParallel 实现",
            "Google ADK 使用 ParallelAgent 实现",
            "需注意：错误处理和结果聚合机制"
        ]
    )
    
    # 保存文件
    output_path = "d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter3_parallelization/Chapter3_Parallelization.pptx"
    prs.save(output_path)
    print(f"Chapter 3 PPT created: {output_path}")

if __name__ == "__main__":
    create_parallelization_ppt()
