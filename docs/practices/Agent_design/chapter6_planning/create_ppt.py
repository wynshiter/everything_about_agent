"""
Chapter 6: Planning PPT Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_planning_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TITLE_COLOR = RGBColor(31, 78, 121)
    
    def add_slide(title, bullets, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.space_after = Pt(12)
    
    def add_title(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            tf = sub.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
    
    add_title("Chapter 6: Planning", "规划模式 - 复杂任务的分解与执行")
    
    add_slide("什么是规划模式？", [
        "将复杂任务分解为可执行的步骤序列",
        "涉及：任务分解、步骤排序、依赖管理",
        "静态规划 vs 动态规划",
        "ReAct: 推理与行动交织"
    ])
    
    add_slide("规划方法", [
        "1. 静态规划：预先定义所有步骤",
        "2. 动态规划：根据上下文实时生成",
        "3. ReAct：Reasoning + Acting 循环"
    ])
    
    add_slide("LangChain 实现", [
        "使用 LLM 生成步骤列表",
        "ChatPromptTemplate 创建规划提示",
        "解析输出为步骤列表",
        "迭代执行每个步骤"
    ])
    
    add_slide("Google ADK 实现", [
        "使用 LoopAgent 实现迭代",
        "Planner Agent 生成计划",
        "Executor Agent 执行步骤",
        "条件检查决定是否继续"
    ])
    
    add_slide("运行指南", [
        "pip install -e .",
        "python src/agents/patterns/planning.py"
    ])
    
    add_slide("总结", [
        "规划是处理复杂任务的核心能力",
        "动态规划更灵活但也更难控制",
        "注意步骤依赖和错误处理"
    ])
    
    prs.save("d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter6_planning/Chapter6_Planning.pptx")
    print("Chapter 6 PPT created")

if __name__ == "__main__":
    create_planning_ppt()
