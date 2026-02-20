"""
Chapter 5: Tool Use PPT Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_tool_use_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    TITLE_COLOR = RGBColor(31, 78, 121)
    
    def add_title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(title: str, bullets: list):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.space_after = Pt(12)
    
    # Slide 1: 封面
    add_title_slide("Chapter 5: Tool Use", "工具使用模式 - 让 Agent 与外部世界交互")
    
    # Slide 2: 什么是工具使用
    add_content_slide("什么是工具使用模式？", [
        "通过 Function Calling 使 Agent 能够调用外部函数",
        "流程：定义工具 -> LLM决策 -> 生成调用 -> 执行 -> 返回结果",
        "工具可以是函数、API、数据库、甚至另一个 Agent",
        "核心：打破 LLM 训练数据的限制，访问实时信息"
    ])
    
    # Slide 3: 工具调用流程
    add_content_slide("工具调用的完整流程", [
        "1. 工具定义：向 LLM 描述可用工具及其参数",
        "2. LLM 决策：判断是否需要调用工具",
        "3. 生成调用：LLM 生成结构化函数调用（JSON）",
        "4. 执行工具：框架执行实际函数",
        "5. 返回结果：工具输出传回 LLM",
        "6. 最终响应：LLM 结合工具结果生成回答"
    ])
    
    # Slide 4: 适用场景
    add_content_slide("工具使用的典型场景", [
        "实时信息：天气、股票、新闻",
        "数据库操作：查询、更新、删除",
        "执行计算：数学运算、数据分析",
        "发送通信：邮件、消息、通知",
        "代码执行：运行代码片段",
        "系统控制：智能家居、IoT设备"
    ])
    
    # Slide 5: LangChain 实现
    add_content_slide("LangChain 实现", [
        "使用 @tool 装饰器定义工具",
        "create_tool_calling_agent() 创建 Agent",
        "AgentExecutor 执行工具调用",
        "bind_tools() 绑定工具到 LLM"
    ])
    
    # Slide 6: Google ADK 工具
    add_content_slide("Google ADK 常用工具", [
        "google_search: 搜索互联网",
        "code_execution: 沙箱中执行 Python",
        "vertex_search: 搜索企业知识库",
        "支持自定义工具扩展"
    ])
    
    # Slide 7: 运行指南
    add_content_slide("运行指南", [
        "pip install langchain-community",
        "定义清晰的工具描述和参数",
        "python src/agents/patterns/tool_use.py"
    ])
    
    # Slide 8: 总结
    add_content_slide("总结与要点", [
        "工具使用是 Agent 区别于普通 LLM 的关键能力",
        "需要清晰的工具定义和错误处理",
        "注意安全性：验证输入、限制权限",
        "避免过度依赖工具，保持核心能力"
    ])
    
    output_path = "d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter5_tool_use/Chapter5_Tool_Use.pptx"
    prs.save(output_path)
    print(f"Chapter 5 PPT created: {output_path}")

if __name__ == "__main__":
    create_tool_use_ppt()
