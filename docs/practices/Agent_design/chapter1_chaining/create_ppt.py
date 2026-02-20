"""
Chapter 1: Prompt Chaining 演示文稿生成脚本
使用 python-pptx 库创建 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 颜色定义 - Ocean Gradient
colors = {
    'primary': RGBColor(0x06, 0x5A, 0x82),      # deep blue
    'secondary': RGBColor(0x1C, 0x72, 0x93),    # teal
    'accent': RGBColor(0x02, 0xC3, 0x9A),       # mint
    'dark': RGBColor(0x21, 0x29, 0x5C),         # midnight
    'light': RGBColor(0xF8, 0xFA, 0xFC),        # off-white
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gray': RGBColor(0x64, 0x74, 0x8B),
    'text_dark': RGBColor(0x1E, 0x29, 0x3B),
    'text_gray': RGBColor(0x47, 0x55, 0x69),
    'code_bg': RGBColor(0x1E, 0x29, 0x3B),
    'code_text': RGBColor(0xE2, 0xE8, 0xF0),
    'code_highlight': RGBColor(0x22, 0xD3, 0xEE),
    'arrow_gray': RGBColor(0x94, 0xA3, 0xB8),
}

def add_text_box(slide, left, top, width, height, text, font_size=12, font_name='Arial', 
                 bold=False, color=None, align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP):
    """添加文本框"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rectangle(slide, left, top, width, height, fill_color):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), 
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ============ Slide 1: 封面 ============
slide_layout = prs.slide_layouts[6]  # 空白布局
slide1 = prs.slides.add_slide(slide_layout)

# 背景
background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
background.fill.solid()
background.fill.fore_color.rgb = colors['dark']
background.line.fill.background()

add_text_box(slide1, 0.5, 1.5, 9, 0.8, "Chapter 1", 24, 'Arial', False, colors['accent'])
add_text_box(slide1, 0.5, 2.2, 9, 1.2, "Prompt Chaining", 48, 'Arial Black', True, colors['white'])
add_text_box(slide1, 0.5, 3.3, 9, 0.6, "提示词链设计模式", 20, 'Arial', False, colors['accent'])

# 装饰线
add_rectangle(slide1, 0.5, 4.2, 2, 0.08, colors['accent'])

add_text_box(slide1, 0.5, 4.5, 9, 0.4, "Everything About Agent - Agent 学习与开发系统", 12, 'Arial', False, colors['arrow_gray'])

# ============ Slide 2: 什么是 Prompt Chaining ============
slide2 = prs.slides.add_slide(slide_layout)

# 背景
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg2.fill.solid()
bg2.fill.fore_color.rgb = colors['light']
bg2.line.fill.background()

# 顶部装饰条
add_rectangle(slide2, 0, 0, 10, 0.15, colors['primary'])

add_text_box(slide2, 0.5, 0.5, 9, 0.7, "什么是 Prompt Chaining?", 32, 'Arial Black', True, colors['dark'])

# 装饰线
add_rectangle(slide2, 0.5, 1.1, 1.5, 0.06, colors['accent'])

# 核心定义卡片
card1 = add_rounded_rectangle(slide2, 0.5, 1.5, 9, 1.3, colors['white'])

definition_text = "Prompt Chaining 是一种将复杂任务分解为一系列较小、较简单的子任务的设计模式。每个子任务由一个独立的 LLM 调用处理，前一个步骤的输出作为下一个步骤的输入。"
add_text_box(slide2, 0.7, 1.7, 8.6, 1, definition_text, 16, 'Arial', False, colors['text_gray'])

# 工作流程标题
add_text_box(slide2, 0.5, 3.1, 9, 0.4, "工作流程", 18, 'Arial', True, colors['dark'])

# 流程图
steps = ["输入", "Step 1", "Step 2", "Step 3", "输出"]
step_colors = [colors['accent'], colors['primary'], colors['secondary'], colors['primary'], colors['accent']]
box_width = 1.6
start_x = 0.6
gap = 0.35

for i, (step, color) in enumerate(zip(steps, step_colors)):
    x = start_x + i * (box_width + gap)
    
    # 方框
    box = add_rounded_rectangle(slide2, x, 3.7, box_width, 0.9, color)
    
    # 文字
    text_box = slide2.shapes.add_textbox(Inches(x), Inches(3.7), Inches(box_width), Inches(0.9))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(14)
    p.font.name = 'Arial'
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_after = Pt(0)
    text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 箭头
    if i < len(steps) - 1:
        arrow_x = x + box_width + 0.05
        arrow = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(4.0), Inches(0.25), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = colors['arrow_gray']
        arrow.line.fill.background()

# ============ Slide 3: 核心优势 ============
slide3 = prs.slides.add_slide(slide_layout)

bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg3.fill.solid()
bg3.fill.fore_color.rgb = colors['light']
bg3.line.fill.background()

add_rectangle(slide3, 0, 0, 10, 0.15, colors['primary'])
add_text_box(slide3, 0.5, 0.5, 9, 0.7, "核心优势", 32, 'Arial Black', True, colors['dark'])
add_rectangle(slide3, 0.5, 1.1, 1.5, 0.06, colors['accent'])

# 三张卡片
advantages = [
    ("降低复杂度", "LLM 在处理单一、专注的任务时表现更好，避免被复杂任务混淆", "🎯"),
    ("提高可靠性", "可以检查和验证每一步的中间结果，确保最终输出质量", "✅"),
    ("便于调试", "容易定位哪个环节出了问题，快速迭代改进", "🔧")
]

for i, (title, desc, icon) in enumerate(advantages):
    x = 0.5 + i * 3.1
    
    # 卡片背景
    add_rounded_rectangle(slide3, x, 1.5, 2.9, 2.8, colors['white'])
    
    # 顶部装饰
    add_rectangle(slide3, x, 1.5, 2.9, 0.1, colors['primary'])
    
    # 图标
    add_text_box(slide3, x, 1.8, 2.9, 0.6, icon, 32, 'Arial', False, None, PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide3, x + 0.15, 2.5, 2.6, 0.5, title, 18, 'Arial', True, colors['dark'], PP_ALIGN.CENTER)
    
    # 描述
    desc_box = slide3.shapes.add_textbox(Inches(x + 0.15), Inches(3.1), Inches(2.6), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.name = 'Arial'
    p.font.color.rgb = colors['gray']
    p.alignment = PP_ALIGN.CENTER

# ============ Slide 4: 练习1 - 基础提取流水线 ============
slide4 = prs.slides.add_slide(slide_layout)

bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg4.fill.solid()
bg4.fill.fore_color.rgb = colors['light']
bg4.line.fill.background()

add_rectangle(slide4, 0, 0, 10, 0.15, colors['secondary'])
add_text_box(slide4, 0.5, 0.5, 9, 0.7, "练习1: 基础提取流水线", 28, 'Arial Black', True, colors['dark'])
add_text_box(slide4, 0.5, 1.1, 9, 0.4, "Basic Extraction Pipeline", 14, 'Arial', False, colors['gray'], PP_ALIGN.LEFT, MSO_ANCHOR.TOP)

# 场景卡片
add_rounded_rectangle(slide4, 0.5, 1.6, 4.3, 1.2, colors['white'])
add_text_box(slide4, 0.7, 1.7, 3.9, 0.35, "场景", 14, 'Arial', True, colors['primary'])
add_text_box(slide4, 0.7, 2.1, 3.9, 0.6, "从非结构化文本中提取技术规格，并转换为标准 JSON 格式", 12, 'Arial', False, colors['text_gray'])

# 流程标题
add_text_box(slide4, 5.2, 1.6, 4.3, 0.4, "处理流程", 14, 'Arial', True, colors['dark'])

# 流程步骤
pipeline1_steps = [
    ("输入", "技术规格文本", colors['accent']),
    ("Step 1", "提取 CPU/内存/存储", colors['primary']),
    ("Step 2", "转换为 JSON", colors['secondary']),
    ("输出", "结构化 JSON", colors['accent'])
]

for i, (title, desc, color) in enumerate(pipeline1_steps):
    y = 2.1 + i * 0.55
    
    # 步骤框
    add_rounded_rectangle(slide4, 5.2, y, 0.9, 0.45, color)
    add_text_box(slide4, 5.2, y, 0.9, 0.45, title, 10, 'Arial', True, colors['white'], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    
    # 描述
    add_text_box(slide4, 6.2, y, 3.3, 0.45, desc, 11, 'Arial', False, colors['text_gray'])

# 代码示例标题
add_text_box(slide4, 0.5, 3.0, 4.3, 0.4, "示例代码结构", 14, 'Arial', True, colors['dark'])

# 代码背景
code_bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.4), Inches(9), Inches(1.6))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = colors['code_bg']
code_bg.line.fill.background()

code_text = '''# 使用 LangChain LCEL 构建链式处理
extraction_chain = prompt_extract | llm | StrOutputParser()
full_chain = {"specifications": extraction_chain} | prompt_transform | llm | StrOutputParser()'''

code_box = slide4.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(8.8), Inches(1.4))
tf = code_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(11)
p.font.name = 'Consolas'
p.font.color.rgb = colors['code_text']

# ============ Slide 5: 练习2 - 创意写作流水线 ============
slide5 = prs.slides.add_slide(slide_layout)

bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg5.fill.solid()
bg5.fill.fore_color.rgb = colors['light']
bg5.line.fill.background()

add_rectangle(slide5, 0, 0, 10, 0.15, colors['secondary'])
add_text_box(slide5, 0.5, 0.5, 9, 0.7, "练习2: 创意写作流水线", 28, 'Arial Black', True, colors['dark'])
add_text_box(slide5, 0.5, 1.1, 9, 0.4, "Creative Writing Pipeline", 14, 'Arial', False, colors['gray'])

# 场景卡片
add_rounded_rectangle(slide5, 0.5, 1.6, 4.3, 1.2, colors['white'])
add_text_box(slide5, 0.7, 1.7, 3.9, 0.35, "场景", 14, 'Arial', True, colors['primary'])
add_text_box(slide5, 0.7, 2.1, 3.9, 0.6, "根据写作主题，逐步生成标题、大纲和引言", 12, 'Arial', False, colors['text_gray'])

# 流程标题
add_text_box(slide5, 5.2, 1.6, 4.3, 0.4, "处理流程", 14, 'Arial', True, colors['dark'])

# 流程步骤
pipeline2_steps = [
    ("输入", "写作主题", colors['accent']),
    ("Step 1", "生成吸引人的标题", colors['primary']),
    ("Step 2", "根据标题生成大纲", colors['secondary']),
    ("Step 3", "撰写引言", colors['primary']),
    ("输出", "完整博客引言", colors['accent'])
]

for i, (title, desc, color) in enumerate(pipeline2_steps):
    y = 2.0 + i * 0.5
    
    add_rounded_rectangle(slide5, 5.2, y, 0.9, 0.4, color)
    add_text_box(slide5, 5.2, y, 0.9, 0.4, title, 9, 'Arial', True, colors['white'], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text_box(slide5, 6.2, y, 3.3, 0.4, desc, 11, 'Arial', False, colors['text_gray'])

# 特点
add_text_box(slide5, 0.5, 3.0, 4.3, 0.4, "特点: 演示顺序依赖关系", 14, 'Arial', True, colors['dark'])

features = "• 每一步依赖前一步的输出\n• 展示多步骤链式处理的典型模式\n• 适合需要逐步构建内容的场景"
add_text_box(slide5, 0.5, 3.4, 4.3, 1.2, features, 12, 'Arial', False, colors['text_gray'])

# ============ Slide 6: 运行指南 ============
slide6 = prs.slides.add_slide(slide_layout)

bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg6.fill.solid()
bg6.fill.fore_color.rgb = colors['light']
bg6.line.fill.background()

add_rectangle(slide6, 0, 0, 10, 0.15, colors['accent'])
add_text_box(slide6, 0.5, 0.5, 9, 0.7, "如何运行", 32, 'Arial Black', True, colors['dark'])

# 方法1卡片
add_rounded_rectangle(slide6, 0.5, 1.4, 4.3, 2.5, colors['white'])
add_rectangle(slide6, 0.5, 1.4, 0.1, 2.5, colors['primary'])

add_text_box(slide6, 0.8, 1.5, 3.8, 0.5, "方法1: 使用启动脚本", 16, 'Arial', True, colors['dark'])
add_text_box(slide6, 0.8, 1.9, 3.8, 0.3, "推荐方式", 11, 'Arial', False, colors['accent'])

method1_text = "1. 找到目录:\n   src/practices/Agent_design/chapter1_chaining\n\n2. 双击运行 run.bat"
add_text_box(slide6, 0.8, 2.3, 3.8, 1.5, method1_text, 11, 'Arial', False, colors['text_gray'])

# 方法2卡片
add_rounded_rectangle(slide6, 5.2, 1.4, 4.3, 2.5, colors['white'])
add_rectangle(slide6, 5.2, 1.4, 0.1, 2.5, colors['secondary'])

add_text_box(slide6, 5.5, 1.5, 3.8, 0.5, "方法2: 命令行运行", 16, 'Arial', True, colors['dark'])

# 命令代码背景
code_bg2 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.1), Inches(3.8), Inches(1.5))
code_bg2.fill.solid()
code_bg2.fill.fore_color.rgb = colors['code_bg']
code_bg2.line.fill.background()

add_text_box(slide6, 5.6, 2.2, 3.6, 1.3, "cd src/practices/Agent_design/chapter1_chaining\npython chapter1_chaining_practice.py", 10, 'Consolas', False, colors['code_highlight'])

# 依赖说明
add_text_box(slide6, 0.5, 4.2, 9, 0.4, "运行环境要求", 14, 'Arial', True, colors['dark'])
add_text_box(slide6, 0.5, 4.6, 9, 0.4, "Python 3.10+  •  langchain  •  langchain-core  •  loguru", 12, 'Arial', False, colors['text_gray'])

# ============ Slide 7: 总结 ============
slide7 = prs.slides.add_slide(slide_layout)

bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
bg7.fill.solid()
bg7.fill.fore_color.rgb = colors['dark']
bg7.line.fill.background()

add_text_box(slide7, 0.5, 0.8, 9, 0.7, "总结", 36, 'Arial Black', True, colors['white'])
add_rectangle(slide7, 0.5, 1.4, 2, 0.06, colors['accent'])

# 关键要点
key_points = [
    "Prompt Chaining 将复杂任务分解为顺序执行的子任务",
    "每个步骤的输出作为下一步的输入",
    "使用 LangChain LCEL 的 | 操作符构建处理链",
    "适合需要逐步处理的场景，如信息提取、内容生成"
]

for i, point in enumerate(key_points):
    # 圆点
    dot = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.8 + i * 0.7), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors['accent']
    dot.line.fill.background()
    
    add_text_box(slide7, 1.0, 1.75 + i * 0.7, 8.5, 0.5, point, 14, 'Arial', False, colors['white'])

# 底部
add_text_box(slide7, 0.5, 4.5, 9, 0.4, "下一章: Routing (路由模式)", 14, 'Arial', False, colors['accent'])

# 保存文件
output_path = r"d:\code\python\everything_about_agent\docs\practices\Agent_design\chapter1_chaining\Chapter1_Prompt_Chaining.pptx"
prs.save(output_path)
print(f"PPT 已保存到: {output_path}")
