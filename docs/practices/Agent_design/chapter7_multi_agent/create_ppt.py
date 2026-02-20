"""
Chapter 7: Multi-Agent Collaboration PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TC = RGBColor(31, 78, 121)
    
    def add_title_slide(t, s):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TC
        p.alignment = PP_ALIGN.CENTER
        if s:
            sb = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
            tf = sb.text_frame
            p = tf.paragraphs[0]
            p.text = s
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(t, bullets):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TC
        cb = sl.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf = cb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(20)
            p.space_after = Pt(12)
    
    add_title_slide("Chapter 7: Multi-Agent Collaboration", "多智能体协作模式")
    add_content_slide("什么是多智能体协作？", [
        "多个专业化 Agent 协同工作解决复杂问题",
        "每个 Agent 有特定角色和目标",
        "通过协作产生超越单个 Agent 的能力",
        "关键：任务分解 + 通信协议"
    ])
    add_content_slide("协作形式", [
        "1. 顺序交接：Agent 间依次传递任务",
        "2. 并行处理：多个 Agent 同时工作",
        "3. 层级结构：Supervisor 管理 Workers",
        "4. 专家团队：不同领域专业协作"
    ])
    add_content_slide("通信结构", [
        "Network: 去中心化点对点",
        "Supervisor: 中央协调器模式",
        "Hierarchical: 多层管理结构",
        "Custom: 自定义混合结构"
    ])
    add_content_slide("CrewAI 实现", [
        "from crewai import Agent, Task, Crew",
        "定义 Role 和 Goal",
        "Task 可设置 context 依赖",
        "Crew.kickoff() 执行"
    ])
    add_content_slide("运行指南", ["pip install -e .", "python src/agents/patterns/multi_agent.py"])
    add_content_slide("总结", ["适合复杂多阶段任务", "提高系统模块化和可扩展性", "注意协调开销和通信成本"])
    
    prs.save("d:/code/python/everything_about_agent/docs/practices/Agent_design/chapter7_multi_agent/Chapter7_Multi_Agent.pptx")
    print("Chapter 7 PPT created")

if __name__ == "__main__":
    create_ppt()
